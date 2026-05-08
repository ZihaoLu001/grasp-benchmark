from __future__ import annotations

from pathlib import Path


SLIDES = [
    (
        "GraspVLA Inner Workings",
        [
            "Goal: explain how GraspVLA works to the benchmark team",
            "Scope: public release plus what we have already run on em14",
        ],
    ),
    (
        "Why This Paper Matters",
        [
            "Synthetic-data pretraining for grasping instead of large real robot logs",
            "Direct end-to-end baseline for our benchmark against modular systems",
            "Public release focuses on deployment: server, simulation, and real-world controller",
        ],
    ),
    (
        "Public Release Pieces",
        [
            "GraspVLA repo: model server and offline test",
            "GraspVLA-playground repo: validation, playground, and LIBERO evaluation",
            "GraspVLA-real-world-controller repo: Franka plus dual-camera client",
            "Full training stack and SynGrasp-1B are not fully public",
        ],
    ),
    (
        "Server API",
        [
            "Inputs: front_view_image, side_view_image, proprio_array, text",
            "Prompt wrapper: In: What action should the robot take to {instruction}?",
            "ZeroMQ service returns action deltas plus debug bbox and pose",
            "serve.py warms up the model before opening the request loop",
        ],
    ),
    (
        "How The Model Produces A Grasp",
        [
            "Token pattern: text_ids, bbox, hist_proprio, cur_proprio, goal, eos",
            "Autoregressive stage predicts bbox and goal tokens",
            "Flow matching stage generates continuous action trajectory",
            "Output is detokenized back to xyz, rpy, and gripper",
        ],
    ),
    (
        "Important Code Paths",
        [
            "prompt.py for the CoT prompt wrapper",
            "token_pattern.py for bbox and goal token layout",
            "model/vla/__init__.py for autoregressive plus flow-matching generation",
            "model/vla/flow_matching.py for the action head",
            "scripts/serve.py for deployment inference",
        ],
    ),
    (
        "Current Lab Status",
        [
            "GraspVLA shared-protocol results are stable for the current simulator suites",
            "Contact-GraspNet shared pipeline completed: 25/90, 20/168, 8/40, and 0/24",
            "Depth+K+segmap proposal-path appendix completed: 40/138",
            "Tracked evidence lives in configs/results/cgn_shared_protocol_h100_20260508.json",
            "AnyGrasp is excluded from current comparative claims until fresh license/runtime revalidation",
        ],
    ),
    (
        "Next Step For The Team",
        [
            "Use README.md and docs/current_benchmark_report.md as the collaborator-facing source of truth",
            "Keep generated slides, plots, and videos under artifacts rather than tracked docs",
            "Report success and speed together because the GraspVLA paper treats speed as a method metric",
            "Use the depth+K+segmap appendix for Contact-GraspNet proposal-path context",
        ],
    ),
]


def main() -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit(
            "python-pptx is required. Install it into the active environment before building slides."
        ) from exc

    output_path = Path(__file__).resolve().parents[1] / "artifacts" / "slides" / "graspvla_inner_workings.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    for index, (title, bullets) in enumerate(SLIDES):
        if index == 0:
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            slide.shapes.title.text = title
            slide.placeholders[1].text = "\n".join(bullets)
            continue

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for bullet_index, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if bullet_index == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(22)

    presentation.save(output_path)
    print(f"Wrote slide deck to {output_path}")


if __name__ == "__main__":
    main()
